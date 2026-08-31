#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
@File    : ppt_parse.py
PPT图文解析工具 - 基于视觉模型对PPT文档进行逐页图文解析
"""
import base64
import json
import logging
from io import BytesIO

from ..client import BigModelClient, BigModelAPIError
from ..utils.file_utils import load_file_bytes, SandboxFileError, FileLoadError

logger = logging.getLogger(__name__)

_MAX_SLIDES = 30


def register_ppt_parse(mcp, client: BigModelClient):
    @mcp.tool()
    async def ppt_multimodal_parse(
        ppt_source: str,
        prompt: str = "请解析这个PPT文档的内容，提取文字和图表信息",
    ) -> str:
        """PPT图文解析工具，对PPT文档进行逐页图文解析。

            适用场景：
            - 用户上传PPT文件并要求提取内容
            - 需要解析PPT中的文字、图表、幻灯片

        Args:
            ppt_source: PPT文件来源。
                        - 优先使用URL地址（如附件的OSS地址/key字段）
                        - 支持upload://引用（通过上传端点获取）
                        - 禁止使用沙箱路径（如/home/ubuntu/...），MCP服务无法访问沙箱文件
                        - 若文件仅在沙箱中，须先用Shell工具上传：curl -F file=@<沙箱路径> http://mcp-multimodal:9100/upload，再使用返回的upload://引用
                        - 支持PPTX格式
            prompt: 对PPT内容的解析要求，默认为提取文字和图表信息

        Returns:
            PPT文档的逐页解析结果文本
        """
        try:
            file_bytes, _ = await load_file_bytes(ppt_source)
        except SandboxFileError as e:
            logger.warning(f"沙箱路径访问: {e.sandbox_path}")
            return json.dumps({"error": str(e)}, ensure_ascii=False)
        except FileLoadError as e:
            logger.error(f"PPT文件加载失败: {e}")
            return json.dumps({"error": f"PPT文件加载失败: {str(e)}"}, ensure_ascii=False)

        try:
            slides_data = _extract_ppt_slides(file_bytes)
            if not slides_data:
                return json.dumps({"error": "PPT解析失败，无法提取幻灯片内容"}, ensure_ascii=False)

            slides_to_process = slides_data[:_MAX_SLIDES]
            all_results = []
            for i, slide_info in enumerate(slides_to_process):
                slide_prompt = f"这是PPT文档第{i + 1}页的内容。{prompt}"

                if slide_info.get("image_base64"):
                    result = await client.vl_chat(
                        prompt=slide_prompt,
                        image_base64_list=[slide_info["image_base64"]],
                    )
                else:
                    text_content = slide_info.get("text", "")
                    if text_content.strip():
                        result = await client.vl_chat(
                            prompt=f"{slide_prompt}\n\n幻灯片文本内容:\n{text_content}",
                        )
                    else:
                        result = f"第{i + 1}页: 无可提取内容"

                all_results.append(f"=== 第{i + 1}页 ===\n{result}")

            if len(slides_data) > _MAX_SLIDES:
                all_results.append(f"\n[注: 文档共{len(slides_data)}页，仅解析前{_MAX_SLIDES}页]")

            return "\n\n".join(all_results)
        except BigModelAPIError as e:
            logger.error(f"PPT解析API失败: {e}")
            return json.dumps({"error": f"PPT解析失败: {str(e)}"}, ensure_ascii=False)
        except Exception as e:
            logger.error(f"PPT解析失败: {e}")
            return json.dumps({"error": f"PPT解析失败: {str(e)}"}, ensure_ascii=False)


def _extract_ppt_slides(ppt_bytes: bytes) -> list:
    """从PPT二进制数据中提取每页的文字和图片"""
    try:
        from pptx import Presentation
        from PIL import Image as PILImage

        prs = Presentation(BytesIO(ppt_bytes))
        slides_data = []

        for slide in prs.slides:
            text_parts = []
            image_b64 = None

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            text_parts.append(text)

                if shape.shape_type == 13:
                    try:
                        img_bytes = shape.image.blob
                        img = PILImage.open(BytesIO(img_bytes))
                        buf = BytesIO()
                        img.save(buf, format="PNG")
                        image_b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
                    except Exception:
                        continue

            slides_data.append({
                "text": "\n".join(text_parts),
                "image_base64": image_b64,
            })

        return slides_data
    except Exception as e:
        logger.error(f"PPT幻灯片提取异常: {e}")
        return []
